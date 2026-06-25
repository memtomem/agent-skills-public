"""Core library for parsing PowerPoint decks (.pptx) into Markdown + JSON.

A slide has no reading order: shapes are positioned freely on a canvas, so the
order python-pptx returns them in (z-order) is not the order a human reads them.
A flat text dump therefore scrambles slides, and it drops what matters most —
tables, charts (whose data is in the XML, recoverable exactly), grouped shapes,
and speaker notes.

So this library, per slide, **orders shapes by geometry** (top, then left — the
slide analogue of reading order), recurses into groups, pulls tables (merged
cells included), charts (categories + series values), and the notes, then emits
the same structured element tree as the other parsers with a confidence score
and flags.

Requires python-pptx. markitdown is an optional cross-check only.
"""

from __future__ import annotations

import json
import os
import re
from dataclasses import dataclass, field, asdict
from typing import Any, Optional

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn


class PptxError(Exception):
    """User-facing error for an unreadable/missing deck."""


@dataclass
class Element:
    type: str                       # heading | paragraph | list_item | table | chart | image | note
    slide: int
    text: str = ""
    level: Optional[int] = None     # list_item: indent depth (0-based)
    rows: Optional[list[list[str]]] = None
    spans: Optional[list[dict]] = None
    chart: Optional[dict] = None
    note: Optional[str] = None

    def to_dict(self) -> dict:
        d = asdict(self)
        # level is part of the list_item contract — keep it there even when 0
        # (0 is a real indent depth, not "absent"); drop it on other types.
        if self.type == "list_item":
            d["level"] = int(d.get("level") or 0)
        else:
            d.pop("level", None)
        return {k: v for k, v in d.items() if v not in (None, [], "")}


CONFIDENCE_VERIFY_BELOW = 0.75


def _pos(shape) -> tuple[int, int]:
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    return int(top), int(left)


def _iter_shapes_in_order(shapes):
    """Flatten groups and yield shapes ordered top-to-bottom, left-to-right."""
    flat = []
    for sh in shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
            flat.extend(_iter_shapes_in_order(sh.shapes))
        else:
            flat.append(sh)
    return sorted(flat, key=_pos)


def _overlap(a, b) -> bool:
    try:
        at, al = _pos(a); bt, bl = _pos(b)
        ar, ab = al + (a.width or 0), at + (a.height or 0)
        br, bb = bl + (b.width or 0), bt + (b.height or 0)
        ix = max(0, min(ar, br) - max(al, bl))
        iy = max(0, min(ab, bb) - max(at, bt))
        return ix > 0 and iy > 0
    except Exception:
        return False


# --------------------------------------------------------------------------
# Tables (merged-cell aware)
# --------------------------------------------------------------------------

def parse_table(graphic_frame) -> tuple[list[list[str]], list[dict]]:
    tbl = graphic_frame.table
    rows: list[list[str]] = []
    spans: list[dict] = []
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)
    for r in range(n_rows):
        row_vals = []
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            if cell.is_spanned:
                row_vals.append("")
                continue
            row_vals.append(re.sub(r"\s+", " ", cell.text).strip())
            if cell.is_merge_origin and (cell.span_height > 1 or cell.span_width > 1):
                spans.append({"row": r, "col": c,
                              "rowspan": int(cell.span_height),
                              "colspan": int(cell.span_width)})
        rows.append(row_vals)
    return rows, spans


# --------------------------------------------------------------------------
# Charts (data is exact — it's in the XML)
# --------------------------------------------------------------------------

def parse_chart(graphic_frame) -> dict:
    chart = graphic_frame.chart
    title = None
    try:
        if chart.has_title:
            title = chart.chart_title.text_frame.text.strip() or None
    except Exception:
        title = None
    cats = []
    series = []
    try:
        cats = [str(c) for c in chart.plots[0].categories]
    except Exception:
        pass
    try:
        for s in chart.series:
            series.append({"name": s.name, "values": [v for v in s.values]})
    except Exception:
        pass
    return {"kind": str(chart.chart_type), "title": title,
            "categories": cats, "series": series}


# --------------------------------------------------------------------------
# Text frames
# --------------------------------------------------------------------------

def _bullet_kind(para) -> Optional[str]:
    """Explicit bullet state from the paragraph's own pPr: 'bullet' | 'none' |
    None (inherit from the placeholder/master). PowerPoint stores bullets as
    a:buChar / a:buAutoNum, and an explicit "no bullet" as a:buNone."""
    pPr = para._p.find(qn("a:pPr"))
    if pPr is None:
        return None
    if pPr.find(qn("a:buNone")) is not None:
        return "none"
    if pPr.find(qn("a:buChar")) is not None or pPr.find(qn("a:buAutoNum")) is not None:
        return "bullet"
    return None


def _is_body_placeholder(sh) -> bool:
    """A text body placeholder (BODY/OBJECT/SUBTITLE) whose paragraphs are
    bulleted by default — as opposed to a title or a free text box."""
    try:
        if not sh.is_placeholder:
            return False
        return sh.placeholder_format.type in (
            PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.SUBTITLE,
        )
    except Exception:
        return False


def _emit_text_frame(tf, slide_idx, is_title, is_body=False) -> list[Element]:
    # The title is captured at slide level (slide["title"]) and shown in the
    # slide header, so we don't re-emit it as a heading element (avoids a
    # duplicate line in the Markdown).
    if is_title:
        return []
    out = []
    for para in tf.paragraphs:
        text = "".join(run.text for run in para.runs).strip() or (para.text or "").strip()
        if not text:
            continue
        level = int(para.level or 0)
        kind = _bullet_kind(para)
        # A level-0 bullet in a body placeholder is still a list item — the most
        # common deck shape — even though PowerPoint leaves its bullet implicit.
        # Honour an explicit buNone, and treat any indented paragraph as a list.
        if kind == "none":
            is_list = False
        elif kind == "bullet":
            is_list = True
        else:
            is_list = is_body or level > 0
        if is_list:
            out.append(Element(type="list_item", slide=slide_idx, text=text, level=level))
        else:
            out.append(Element(type="paragraph", slide=slide_idx, text=text))
    return out


# --------------------------------------------------------------------------
# Assembly
# --------------------------------------------------------------------------

def parse_slide(slide, idx: int) -> dict:
    elements: list[Element] = []
    flags: list[str] = []
    confidence = 1.0

    title_shape = slide.shapes.title
    ordered = _iter_shapes_in_order(slide.shapes)

    grouped = any(sh.shape_type == MSO_SHAPE_TYPE.GROUP for sh in slide.shapes)
    n_text_shapes = 0
    n_visual = 0

    for sh in ordered:
        try:
            if sh.has_table:
                rows, spans = parse_table(sh)
                if any(any(r) for r in rows):
                    el = Element(type="table", slide=idx, rows=rows)
                    if spans:
                        el.spans = spans
                        if "table-merged-cells" not in flags:
                            flags.append("table-merged-cells")
                    elements.append(el)
                continue
            if sh.has_chart:
                ch = parse_chart(sh)
                elements.append(Element(type="chart", slide=idx,
                                        text=ch.get("title") or ch.get("kind"), chart=ch))
                if "has-chart" not in flags:
                    flags.append("has-chart")
                continue
        except Exception:
            pass
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            n_visual += 1
            elements.append(Element(type="image", slide=idx,
                                    text=(sh.name or "image"),
                                    note="picture present; not extracted in this pass"))
            continue
        if sh.has_text_frame and sh.text_frame.text.strip():
            n_text_shapes += 1
            is_title = (title_shape is not None and sh == title_shape)
            is_body = _is_body_placeholder(sh)
            elements.extend(_emit_text_frame(sh.text_frame, idx, is_title, is_body))

    # speaker notes
    notes = ""
    try:
        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
    except Exception:
        notes = ""
    if notes:
        elements.append(Element(type="note", slide=idx, text=notes, note="speaker-notes"))

    # confidence flags
    text_shapes = [s for s in ordered if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()]
    overlaps = any(_overlap(a, b) for i, a in enumerate(text_shapes) for b in text_shapes[i + 1:])
    if overlaps:
        confidence = min(confidence, 0.70)
        flags.append("overlapping-shapes-order-ambiguous")
    if grouped:
        confidence = min(confidence, 0.85)
        flags.append("grouped-shapes")
    # An image-only slide carries content local tools can't read — mark it for a
    # vision pass explicitly (parity with pdf-parser's needs_vision), not just a
    # low score, so callers can route it to OCR/transcription deterministically.
    needs_vision = (n_text_shapes == 0 and n_visual > 0)
    if needs_vision:
        confidence = min(confidence, 0.40)
        flags.append("image-only-slide")
        flags.append("needs-vision")
    if "table-merged-cells" in flags:
        confidence = min(confidence, 0.80)
    if not elements:
        confidence = min(confidence, 0.50)
        flags.append("empty-slide")

    title_text = title_shape.text.strip() if (title_shape is not None and title_shape.has_text_frame) else ""
    return {
        "index": idx,
        "title": title_text,
        "confidence": round(confidence, 2),
        "needs_vision": needs_vision,
        "flags": flags,
        "elements": [e.to_dict() for e in elements],
    }


def build_document(path: str) -> dict:
    if not os.path.exists(path):
        raise PptxError(f"file not found: {path!r}")
    try:
        prs = Presentation(path)
    except Exception as e:
        raise PptxError(f"cannot open {path!r}: {e}") from e
    slides = [parse_slide(s, i + 1) for i, s in enumerate(prs.slides)]
    min_conf = round(min((s["confidence"] for s in slides), default=1.0), 2)
    return {
        "source": os.path.basename(path),
        "slide_count": len(slides),
        "min_confidence": min_conf,
        "verify_slides": [s["index"] for s in slides if s["confidence"] < CONFIDENCE_VERIFY_BELOW],
        "vision_slides": [s["index"] for s in slides if s.get("needs_vision")],
        "slides": slides,
    }


# --------------------------------------------------------------------------
# Rendering
# --------------------------------------------------------------------------

def _esc(s):
    return (s or "").replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _rows_to_gfm(rows):
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    norm = [r + [""] * (width - len(r)) for r in rows]
    def fmt(r):
        # escape HTML-sensitive chars too — a GFM cell is still rendered to HTML.
        return "| " + " | ".join(_esc(c).replace("|", "\\|").replace("\n", " ") for c in r) + " |"
    out = [fmt(norm[0]), "| " + " | ".join(["---"] * width) + " |"]
    out += [fmt(r) for r in norm[1:]]
    return "\n".join(out)


def _rows_to_html(rows, spans):
    spans = spans or []
    anchor = {(s["row"], s["col"]): (s.get("rowspan", 1), s.get("colspan", 1)) for s in spans}
    covered = set()
    for s in spans:
        for dr in range(s.get("rowspan", 1)):
            for dc in range(s.get("colspan", 1)):
                if dr or dc:
                    covered.add((s["row"] + dr, s["col"] + dc))
    out = ["<table>"]
    for ri, row in enumerate(rows):
        tag = "th" if ri == 0 else "td"
        cells = []
        for ci, val in enumerate(row):
            if (ri, ci) in covered:
                continue
            attr = ""
            if (ri, ci) in anchor:
                rs, cs = anchor[(ri, ci)]
                if rs > 1:
                    attr += f' rowspan="{rs}"'
                if cs > 1:
                    attr += f' colspan="{cs}"'
            cells.append(f"<{tag}{attr}>{_esc(val).replace(chr(10), ' ')}</{tag}>")
        out.append("  <tr>" + "".join(cells) + "</tr>")
    out.append("</table>")
    return "\n".join(out)


def element_to_markdown(el: dict) -> str:
    # User-authored text is HTML-escaped on the way into Markdown (it may be
    # rendered back to HTML downstream); JSON keeps the raw values.
    t = el.get("type")
    text = _esc((el.get("text") or "").strip())
    if t == "heading":
        return f"## {text}"
    if t == "paragraph":
        return text
    if t == "list_item":
        indent = "  " * int(el.get("level") or 0)
        return f"{indent}- {text}"
    if t == "note":
        return f"> 🗒 **[notes]** {text}"
    if t == "image":
        return f"<!-- 🖼 {text}: {_esc(el.get('note', ''))} -->"
    if t == "chart":
        ch = el.get("chart") or {}
        lines = [f"**[chart]** {text} _(kind: {_esc(str(ch.get('kind', '?')))})_"]
        cats = ch.get("categories") or []
        for s in ch.get("series") or []:
            pairs = ", ".join(f"{_esc(str(c))}={_esc(str(v))}"
                              for c, v in zip(cats, s.get("values") or []))
            lines.append(f"- {_esc(str(s.get('name')))}: {pairs}")
        return "\n".join(lines)
    if t == "table":
        rows = el.get("rows") or []
        spans = el.get("spans")
        return _rows_to_html(rows, spans) if spans else _rows_to_gfm(rows)
    return el.get("text", "").strip()


def to_markdown(document: dict) -> str:
    out = [f"<!-- parsed from {_esc(document.get('source', ''))} -->\n"]
    for sl in document.get("slides", []):
        head = _esc(sl.get("title") or f"Slide {sl['index']}")
        out.append(f"# Slide {sl['index']}: {head}")
        conf = sl.get("confidence")
        if conf is not None and conf < CONFIDENCE_VERIFY_BELOW:
            out.append(f"> 🔎 Slide {sl['index']} low parse confidence ({conf}); "
                       f"verify. flags: {', '.join(sl.get('flags') or []) or 'none'}")
        for el in sl.get("elements", []):
            md = element_to_markdown(el)
            if md:
                out.append(md)
        out.append("")
    return "\n\n".join(out).strip() + "\n"


def save_outputs(document: dict, md_path: str, json_path: str) -> None:
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(document, f, ensure_ascii=False, indent=2)
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(to_markdown(document))


def markitdown_crosscheck(path: str) -> Optional[str]:
    try:
        from markitdown import MarkItDown
    except Exception:
        return None
    try:
        return MarkItDown().convert(path).text_content
    except Exception:
        return None
