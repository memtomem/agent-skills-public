"""Generate a messy .pptx fixture (never committed; built on demand).

deck.pptx:
  slide 1: title + bullet list + speaker notes
  slide 2: a table with a merged header cell (colspan) + a bar chart (data)
  slide 3: an image-only slide (picture, no text)
"""
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from PIL import Image


def _slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. 개요 Overview"
    body = slide.placeholders[1].text_frame
    body.text = "디지털 전환 가속"
    for t in ["해외 시장 확대", "R&D 투자 증대"]:
        p = body.add_paragraph(); p.text = t; p.level = 1
    slide.notes_slide.notes_text_frame.text = "발표자 노트: 3분 분량으로 설명할 것."


def _slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "2. 재무 Financials"
    # table with merged header
    rows, cols = 3, 3
    gf = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(5), Inches(2))
    tbl = gf.table
    tbl.cell(0, 0).text = "Region"
    tbl.cell(0, 1).merge(tbl.cell(0, 2))
    tbl.cell(0, 1).text = "Sales H1"
    data = [["Asia", "120", "135"], ["EU", "90", "110"]]
    for i, row in enumerate(data):
        for j, v in enumerate(row):
            tbl.cell(i + 1, j).text = v
    # bar chart with real data
    cd = CategoryChartData()
    cd.categories = ["Q1", "Q2", "Q3", "Q4"]
    cd.add_series("Revenue", (260, 295, 330, 354))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                           Inches(6), Inches(1.5), Inches(3.2), Inches(3), cd)


def _slide3(prs, tmp_png):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    img = Image.new("RGB", (400, 300), "navy")
    img.save(tmp_png)
    slide.shapes.add_picture(tmp_png, Inches(1), Inches(1), Inches(4), Inches(3))


def build(path: str) -> str:
    prs = Presentation()
    _slide1(prs)
    _slide2(prs)
    import tempfile
    fd, tmp = tempfile.mkstemp(suffix=".png"); os.close(fd)
    _slide3(prs, tmp)
    prs.save(path)
    try:
        os.remove(tmp)
    except OSError:
        pass
    return path


def make_all(out_dir: str) -> None:
    os.makedirs(out_dir, exist_ok=True)
    build(os.path.join(out_dir, "deck.pptx"))


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else os.path.dirname(__file__)
    make_all(out)
    print("fixtures written to", out)
